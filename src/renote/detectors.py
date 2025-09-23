from __future__ import annotations

from typing import Optional

from pptx.enum.shapes import MSO_SHAPE_TYPE


def _iter_shapes(slide) -> list:
    """Flatten shapes including groups.

    Returns a list of shapes, recursing into group shapes.
    """
    shapes = list(slide.shapes)
    flat = []
    while shapes:
        shape = shapes.pop()
        if getattr(shape, "shape_type", None) == MSO_SHAPE_TYPE.GROUP or hasattr(shape, "shapes"):
            try:
                shapes.extend(shape.shapes)
            except Exception:
                pass
            continue
        flat.append(shape)
    return flat


def get_shape_font_sizes(shape) -> Optional[list[float]]:
    """Return list of font sizes (pt) for a text shape; None if no text frame.

    If run/paragraph sizes are not explicitly set (inherit theme), fall back to 18pt.
    """
    if not getattr(shape, "has_text_frame", False):
        return None
    sizes: list[float] = []
    tf = shape.text_frame
    for p in tf.paragraphs:
        para_size = None
        try:
            if p.font is not None and p.font.size is not None:
                para_size = p.font.size.pt
        except Exception:
            para_size = None
        for r in p.runs:
            try:
                if r.font.size is not None:
                    sizes.append(r.font.size.pt)
                else:
                    sizes.append(para_size or 18.0)
            except Exception:
                sizes.append(para_size or 18.0)
    return sizes or [18.0]


def is_title_slide(slide, third_width: int, *, min_font_pt: float = 36.0, min_width_ratio: float = 1.2) -> bool:
    """Heuristic: slide is title-like if it has a wide text shape with large font.

    Args:
        slide: python-pptx slide
        third_width: reference width (one third of original)
        min_font_pt: minimal font size to qualify as title
        min_width_ratio: minimal shape width ratio vs third_width
    """
    for shape in _iter_shapes(slide):
        try:
            shape_width = int(shape.width)
        except Exception:
            continue
        if shape_width >= int(min_width_ratio * third_width) and getattr(shape, "has_text_frame", False):
            sizes = get_shape_font_sizes(shape)
            if sizes and any(size >= min_font_pt for size in sizes):
                return True
    return False


def pick_primary_title_shape(slide) -> Optional[object]:
    """Pick main title text shape by largest median font size then width.

    If sizes are not explicit, use 18pt fallback so that a candidate is still chosen.
    """
    candidates = []
    for shape in _iter_shapes(slide):
        if not getattr(shape, "has_text_frame", False):
            continue
        sizes = get_shape_font_sizes(shape) or [18.0]
        sorted_sizes = sorted(sizes)
        median = sorted_sizes[len(sorted_sizes) // 2]
        candidates.append((median, int(getattr(shape, "width", 0)), shape))
    if not candidates:
        return None
    candidates.sort(key=lambda x: (x[0], x[1]), reverse=True)
    return candidates[0][2]


