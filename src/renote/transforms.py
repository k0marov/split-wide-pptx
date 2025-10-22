from typing import Iterable, Optional, List
from io import BytesIO

from pptx import Presentation
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.util import Pt
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_SHAPE_TYPE, MSO_AUTO_SHAPE_TYPE
from pptx.oxml.ns import qn

from .detectors import is_title_slide, pick_primary_title_shape

def _clone_shape(slide, src_slide, shape):
    from copy import deepcopy
    try:
        new_el = deepcopy(shape._element)
    except Exception:
        return
    _fix_relationships_for_element(new_el, src_slide, slide)
    try:
        slide.shapes._spTree.insert_element_before(new_el, "p:extLst")
    except Exception:
        try:
            slide.shapes._spTree.append(new_el)
        except Exception:
            return
    return new_el


def _delete_shape(shape) -> None:
    element = shape._element
    parent = element.getparent()
    if parent is not None:
        parent.remove(element)


def _iter_shapes(slide):
    for shape in list(slide.shapes):
        yield shape

def _cut_layouts(prs):
    """Cuts layouts so that the background does not overflow the new slide width."""
    target_width = prs.slide_width
    for master_idx, slide_master in enumerate(prs.slide_masters):
        for shape in slide_master.shapes:
            if shape.left is None:
                continue
            if shape.width > target_width:
                shape.width = target_width
                if shape.left == 0:
                    shape.left = 0
        for layout_idx, layout in enumerate(slide_master.slide_layouts):
            for shape in layout.shapes:
                if shape.left is None:
                    continue
                if shape.left + shape.width > target_width:
                    new_width = target_width - shape.left
                    if new_width > 0:
                        shape.width = new_width
                    # else:
                    #     print("hard case")

    # for slide in prs.slides:
    #     for shape in slide.shapes:
    #         # Resize shapes that extend beyond new width
    #         if shape.left + shape.width > target_width or shape.left < 0:
    #             _delete_shape(shape)


def _remove_placeholders(slide) -> None:
    """Remove placeholders from a slide to avoid duplicates on cloning."""
    for shp in list(slide.shapes):
        if getattr(shp, "is_placeholder", False):
            _delete_shape(shp)


def _fix_relationships_for_element(new_el, src_slide, dst_slide) -> None:
    """Rebind relationship ids for copied XML subtree to destination slide."""
    REL_NS = "http://schemas.openxmlformats.org/officeDocument/2006/relationships"
    ATTRS = (f"{{{REL_NS}}}embed", f"{{{REL_NS}}}id")

    old_to_new = {}
    try:
        for old_rid, rel in src_slide.part.rels.items():
            try:
                new_rid = dst_slide.part.relate_to(rel._target, rel.reltype)
                old_to_new[old_rid] = new_rid
            except Exception:
                continue
    except Exception:
        pass

    for el in new_el.iter():
        for attr_name in ATTRS:
            old_rid = el.get(attr_name)
            if not old_rid:
                continue
            new_rid = old_to_new.get(old_rid)
            if new_rid:
                el.set(attr_name, new_rid)


def _copy_background(src_slide, dst_slide) -> None:
    """Copy custom background (cSld.bg) and fix relationships."""
    try:
        src_cSld = src_slide._element.cSld
        dst_cSld = dst_slide._element.cSld
    except Exception:
        return
    src_bg = getattr(src_cSld, "bg", None)
    if src_bg is None:
        return
    try:
        dst_bg = getattr(dst_cSld, "bg", None)
        if dst_bg is not None:
            dst_cSld.remove(dst_bg)
    except Exception:
        pass
    from copy import deepcopy
    new_bg = deepcopy(src_bg)
    _fix_relationships_for_element(new_bg, src_slide, dst_slide)
    try:
        dst_cSld.insert(0, new_bg)
        try:
            dst_slide.follow_master_background = False
        except Exception:
            pass
    except Exception:
        try:
            dst_cSld.append(new_bg)
        except Exception:
            pass


def _delete_slide(prs: Presentation, index: int) -> None:
    """Delete slide by dropping relationship and removing from sldIdLst."""
    sldIdLst = prs.slides._sldIdLst
    sldId = sldIdLst[index]
    try:
        rId = sldId.rId  # type: ignore[attr-defined]
    except Exception:
        rId = sldId.get(qn("r:id"))
    if rId:
        try:
            prs.part.drop_rel(rId)
        except Exception:
            pass
    sldIdLst.remove(sldId)


def transform_title_slide(slide, third_width: int) -> None:
    """Reduce slide to a single cloned title shape preserving styles; adjust geometry only."""
    primary = pick_primary_title_shape(slide)
    if primary is None:
        for shp in slide.shapes:
            if getattr(shp, "has_text_frame", False):
                primary = shp
                break
    if primary is None:
        return

    # Remove all shapes on slide
    for shp in list(slide.shapes):
        _delete_shape(shp)

    # Clone original text shape element to preserve all formatting/themes
    _clone_shape(slide, primary)

    # Adjust geometry to fit third width and center vertically
    try:
        shp = slide.shapes[-1]
        shp.width = third_width
        shp.left = 0
        try:
            shp.height = slide.part.presentation.slide_height  # type: ignore[attr-defined]
            shp.top = 0
        except Exception:
            pass
        if getattr(shp, "has_text_frame", False):
            tf = shp.text_frame
            # Keep horizontal alignment as-is if set; otherwise center
            try:
                for p in tf.paragraphs:
                    p.alignment = p.alignment or PP_ALIGN.CENTER
            except Exception:
                for p in tf.paragraphs:
                    p.alignment = PP_ALIGN.CENTER
            try:
                tf.vertical_anchor = MSO_ANCHOR.MIDDLE
            except Exception:
                pass
    except Exception:
        pass




class InvalidSlideCountException(Exception):
    pass

def create_triptych(
    input_pptx: str,
    output_pptx: str,
    *,
    title_min_font_pt: float = 36.0,
    title_min_width_ratio: float = 1.2,
    scenarios: Optional[List[str]] = None,
) -> None:
    """Direct mode preserving styles: duplicate each slide three times in-place then filter per third.

    Мы сохраняем тему, фон и стили, создавая новые слайды в рамках исходной презентации,
    а затем удаляем оригинальные. Экспорт в output_pptx уже с узкой шириной.
    """
    prs = Presentation(input_pptx)
    if len(prs.slides) % 3 != 0:
        raise InvalidSlideCountException()

    original_width = int(prs.slide_width)
    big_width = int(original_width * 3)

    original_count = len(prs.slides)
    originals = [prs.slides[i] for i in range(original_count)]

    from copy import deepcopy
    for src_index in range(0, len(originals), 3):
        main_src = originals[src_index]
        second_src = originals[src_index+1]
        third_src = originals[src_index+2]

        new_slide = prs.slides.add_slide(main_src.slide_layout)
        _remove_placeholders(new_slide)
        for offset, src in ((0, main_src), (original_width, second_src), (2*original_width, third_src)):
            # copy shapes xml
            for shape in src.shapes:
                shape.left += offset
                try:
                    new_el = deepcopy(shape._element)
                except Exception:
                    continue
                _fix_relationships_for_element(new_el, main_src, new_slide)
                try:
                    new_slide.shapes._spTree.insert_element_before(new_el, "p:extLst")
                except Exception:
                    try:
                        new_slide.shapes._spTree.append(new_el)
                    except Exception:
                        continue
            # copy background
            _copy_background(main_src, new_slide)
    # # 2) Удалить оригинальные слайды (они первые original_count в последовательности)
    # # for idx in range(original_count - 1, -1, -1):
    # #     _delete_slide(prs, idx)
    #
    # # 3) Для каждого трио отфильтровать и сдвинуть элементы под свою треть
    # prs.slide_width = big_width
    # for i in range(len(prs.slides)-1, -1, -3):
    #     main_slide = prs.slides[i]
    #     second_slide = prs.slides[i-1]
    #     third_slide = prs.slides[i-2]
    #     _remove_placeholders(main_slide)
    #     for offset, copy_slide in ((original_width, second_slide), (original_width*2, third_slide)):
    #         for shape in list(copy_slide.shapes):
    #             try:
    #                 shape_left = int(shape.left)
    #                 shape_width = int(shape.width)
    #             except Exception:
    #                 _delete_shape(shape)
    #                 continue
    #             try:
    #                 new_shape = _clone_shape(main_slide, shape)
    #                 new_shape.left += offset
    #             except:
    #                 continue
    #
    prs.save(output_pptx)


