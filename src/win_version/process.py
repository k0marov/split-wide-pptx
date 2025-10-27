import json

from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.enum.shapes import MSO_SHAPE_TYPE
import sys
from pptx import Presentation
from pptx.util import Pt

def _delete_shape(shape) -> None:
    element = shape._element
    parent = element.getparent()
    if parent is not None:
        parent.remove(element)

def split_slides_into_thirds(input_pptx, slide_mapping: dict, output_pptx):
    # Load presentation
    prs = Presentation(input_pptx)

    # Calculate dimensions
    original_width = prs.slide_width
    third_width = original_width // 3
    prs.slide_width = third_width

    for i, slide in enumerate(prs.slides):
        # Set boundaries for each third
        info = slide_mapping[str(i)]
        if info['clone_index'] == 1:  # First third
            left_bound, right_bound = 0, third_width
        elif info['clone_index'] == 2:  #
            left_bound, right_bound = third_width, 2 * third_width
        else:  # Last third
            left_bound, right_bound = 2 * third_width, original_width

        is_title = info['type'] == 'title'
        print(is_title) 
        shapes = list(slide.shapes)
        max_left_delta = 0
        while shapes:
            shape = shapes.pop()

            shape_left = int(shape.left)
            shape_width = int(shape.width)
            shape_right = shape_left + shape_width

            spans_all_thirds = (shape_left < third_width and
                               shape_right > 2 * third_width)
            if is_title:
                if getattr(shape, 'text', None) is None or i % 3 != 0: 
                    sp = shape._element
                    sp.getparent().remove(sp)
                else: 
                    shape.width = third_width 
                    shape.left = 0 
                    for p in shape.text_frame.paragraphs: 
                        p.alignment = PP_ALIGN.CENTER

                continue 

            if not is_title and (shape.left > right_bound or shape.left+shape.width < left_bound):
                _delete_shape(shape)
                continue
            else:
                if info['clone_index'] == 2:  # Middle third
                    shape.left = shape.left-third_width
                elif info['clone_index'] == 3:  # Last third
                    shape.left = shape.left-2*third_width
                if shape.left < 0:
                    max_left_delta = max(max_left_delta, abs(shape.left))
        shapes = list(slide.shapes)
        while shapes:
            shape = shapes.pop()
            shape.left += max_left_delta


    prs.save(output_pptx)
    print(f"Presentation saved to {output_pptx}")

# if __name__ == "__main__":
#     if len(sys.argv) < 2:
#         print("Usage: python script.py input.pptx [output.pptx]")
#         sys.exit(1)
#
#     input_file = sys.argv[1]
#     output_file = sys.argv[2] if len(sys.argv) > 2 else "split_presentation.pptx"
#
#     split_slides_into_thirds(input_file, output_file)
