from __future__ import annotations

import base64
import os
import shutil
import subprocess
import tempfile
from typing import List, Optional

import requests
from langchain_ollama import ChatOllama
from langchain_core.messages import HumanMessage
from pptx import Presentation

from .detectors import is_title_slide


def _which_soffice(soffice_path: Optional[str] = None) -> Optional[str]:
    if soffice_path and os.path.isfile(soffice_path):
        return soffice_path
    return shutil.which("soffice")


def export_slides_to_pngs(pptx_path: str, out_dir: str, soffice_path: Optional[str] = None) -> List[str]:
    soffice = _which_soffice(soffice_path)
    if not soffice:
        if soffice_path:
            raise RuntimeError(f"LibreOffice не найден по указанному пути: {soffice_path}. Проверьте правильность пути.")
        else:
            raise RuntimeError("LibreOffice (soffice) не найден в системе. Установите LibreOffice и укажите путь к soffice.")
    
    cmd = [
        soffice,
        "--headless",
        "--convert-to",
        "png",
        "--outdir",
        out_dir,
        pptx_path,
    ]
    
    try:
        result = subprocess.run(cmd, check=True, stdout=subprocess.PIPE, stderr=subprocess.PIPE)
    except subprocess.CalledProcessError as e:
        stderr_output = e.stderr.decode('utf-8') if e.stderr else "Нет информации об ошибке"
        raise RuntimeError(f"Ошибка при экспорте слайдов в PNG через LibreOffice: {stderr_output}")
    except FileNotFoundError:
        raise RuntimeError(f"Не удается запустить LibreOffice по пути: {soffice}")
    
    pngs = [os.path.join(out_dir, f) for f in os.listdir(out_dir) if f.lower().endswith(".png")]
    if not pngs:
        raise RuntimeError("LibreOffice не создал PNG файлы. Возможно, презентация повреждена или формат не поддерживается.")
    
    pngs.sort()
    return pngs


def _classify_image_langchain(image_path: str, model: str, url: str, timeout: int = 60) -> str:
    try:
        with open(image_path, "rb") as f:
            image_b64 = base64.b64encode(f.read()).decode("utf-8")
    except FileNotFoundError:
        raise RuntimeError(f"Не найден PNG файл для анализа: {image_path}")
    except Exception as e:
        raise RuntimeError(f"Ошибка при чтении PNG файла {image_path}: {str(e)}")
    
    prompt = (
        "Классифицируй слайд презентации. Ответь одним словом: 'title' если это титульный/заголовочный слайд"
        " с крупным заголовком и минимумом контента; иначе ответь 'split'. Без пояснений."
    )
    
    try:
        llm = ChatOllama(model=model, base_url=url, temperature=0)
        msg = HumanMessage(
            content=[
                {"type": "text", "text": prompt},
                {"type": "image_url", "image_url": f"data:image/png;base64,{image_b64}"},
            ]
        )
        out = llm.invoke([msg])
    except Exception as e:
        error_msg = str(e).lower()
        if "connection" in error_msg or "refused" in error_msg:
            raise RuntimeError(f"Не удается подключиться к Ollama по адресу {url}. Убедитесь, что запущен 'ollama serve'.")
        elif "model" in error_msg or "not found" in error_msg:
            raise RuntimeError(f"Модель '{model}' не найдена в Ollama. Выполните: ollama pull {model}")
        else:
            raise RuntimeError(f"Ошибка при обращении к Ollama: {str(e)}")
    
    text = (out.content or "").strip().lower()
    if "title" in text and "split" not in text:
        return "title"
    if text.startswith("title"):
        return "title"
    return "split"


def classify_slides(
    pptx_path: str,
    *,
    mode: str = "heuristic",
    title_min_font_pt: float = 36.0,
    title_min_width_ratio: float = 1.2,
    soffice_path: Optional[str] = None,
    ollama_model: str = "gemma3:4b",
    ollama_url: str = "http://localhost:11434",
) -> List[str]:
    if mode not in {"heuristic", "ollama"}:
        mode = "heuristic"

    if mode == "heuristic":
        prs = Presentation(pptx_path)
        third_width = int(prs.slide_width // 3)
        return [
            "title" if is_title_slide(slide, third_width, min_font_pt=title_min_font_pt, min_width_ratio=title_min_width_ratio) else "split"
            for slide in prs.slides
        ]

    with tempfile.TemporaryDirectory() as tmp:
        images = export_slides_to_pngs(pptx_path, tmp, soffice_path=soffice_path)
        prs = Presentation(pptx_path)
        if len(images) != len(prs.slides):
            images = sorted(images, key=lambda p: os.path.getctime(p))
        return [_classify_image_langchain(img, ollama_model, ollama_url) for img in images]


