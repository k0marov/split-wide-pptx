from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.enum.shapes import MSO_SHAPE_TYPE
import sys
from pptx import Presentation
from pptx.util import Pt

def split_slides_into_thirds(input_pptx, output_pptx):
    # Load presentation
    prs = Presentation(input_pptx)

    # Calculate dimensions
    original_width = prs.slide_width
    third_width = original_width // 3
    prs.slide_width = third_width

    for i, slide in enumerate(prs.slides):
        # Set boundaries for each third
        if i % 3 == 0:  # First third
            left_bound, right_bound = 0, third_width
        elif i % 3 == 1:  # Middle third
            left_bound, right_bound = third_width, 2 * third_width
        else:  # Last third
            left_bound, right_bound = 2 * third_width, original_width

        is_title = check_is_title(slide, third_width)
        print(is_title) 
        shapes = list(slide.shapes)
        while shapes:
            shape = shapes.pop()
            #if shape.shape_type == MSO_SHAPE_TYPE.GROUP: 
            #    shapes += shape.shapes
            #    continue
            #if hasattr(shape, 'shapes'): 
            #    shapes += shape.shapes

            shape_left = int(shape.left)
            shape_width = int(shape.width)
            shape_right = shape_left + shape_width

            spans_all_thirds = (shape_left < third_width and
                               shape_right > 2 * third_width)
            if is_title:
                if getattr(shape, 'text', None) is None or i % 3 != 0: 
                    sp = shape._element
                    sp.getparent().remove(sp)
                else: 
                    shape.width = third_width 
                    shape.left = 0 
                    for p in shape.text_frame.paragraphs: 
                        p.alignment = PP_ALIGN.CENTER

                continue 

            if not is_title and (shape.left > right_bound or shape.left+shape.width < left_bound):
                sp = shape._element
                sp.getparent().remove(sp)
            else:
                if i % 3 == 1:  # Middle third
                    shape.left -= third_width
                elif i % 3 == 2:  # Last third
                    shape.left -= 2 * third_width

    prs.save(output_pptx)
    print(f"Presentation saved to {output_pptx}")

TITLE_FONT_SIZE = 40 

def get_shape_font_size(shape):
    """Returns list of font sizes used in shape's text"""
    font_sizes = []

    if not shape.has_text_frame:
        return None  # Shape contains no text

    for paragraph in shape.text_frame.paragraphs:
        for run in paragraph.runs:
            if run.font.size:
                font_sizes.append(run.font.size.pt)
            else:
                # When no explicit size, PowerPoint uses default (usually 18pt for title, 12pt for body)
                font_sizes.append(12)
    print(font_sizes)
    return font_sizes

def check_is_title(slide, third_width): 
    # TODO: maybe use ML for this 
    shapes = list(slide.shapes)
    while shapes:
        shape = shapes.pop()
        if hasattr(shape, 'shapes'): 
            shapes += shape.shapes
        if shape.shape_type == MSO_SHAPE_TYPE.GROUP: 
            shapes += shape.shapes
            continue

        shape_left = int(shape.left)
        shape_width = int(shape.width)
        shape_right = shape_left + shape_width

        spans_all_thirds = shape_width >= 1.3*third_width 

        if spans_all_thirds and shape.has_text_frame and any(map( lambda fs: fs >= TITLE_FONT_SIZE, get_shape_font_size(shape),)):
            return True
    return False 



if __name__ == "__main__":
    if len(sys.argv) < 2:
        print("Usage: python script.py input.pptx [output.pptx]")
        sys.exit(1)

    input_file = sys.argv[1]
    output_file = sys.argv[2] if len(sys.argv) > 2 else "split_presentation.pptx"

    split_slides_into_thirds(input_file, output_file)
